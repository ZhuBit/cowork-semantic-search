"""Per-format text extraction from document files."""

import csv
import io
from pathlib import Path

SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx", ".csv"}


def extract_text(file_path: Path) -> list[dict]:
    """Extract text from a file, returning list of {text, metadata} dicts.

    Metadata may include page_number (PDF), slide_number (PPTX).
    """
    suffix = file_path.suffix.lower()

    match suffix:
        case ".txt" | ".md":
            text = file_path.read_text(encoding="utf-8", errors="replace")
            return [{"text": text, "metadata": {}}]
        case ".pdf":
            return _extract_pdf(file_path)
        case ".docx":
            return _extract_docx(file_path)
        case ".pptx":
            return _extract_pptx(file_path)
        case ".csv":
            return _extract_csv(file_path)
        case _:
            raise ValueError(
                f"Unsupported file type: {suffix}. "
                f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
            )


def _extract_pdf(file_path: Path) -> list[dict]:
    import pymupdf

    doc = pymupdf.open(str(file_path))
    parts = []
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text()
        parts.append({"text": text.strip(), "metadata": {"page_number": page_num}})
    doc.close()
    return parts


def _extract_docx(file_path: Path) -> list[dict]:
    from docx import Document

    doc = Document(str(file_path))
    text = "\n\n".join(p.text for p in doc.paragraphs)
    return [{"text": text, "metadata": {}}]


def _extract_pptx(file_path: Path) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(str(file_path))
    parts = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        parts.append({
            "text": "\n".join(texts),
            "metadata": {"slide_number": slide_num},
        })
    return parts


def _extract_csv(file_path: Path) -> list[dict]:
    raw = file_path.read_text(encoding="utf-8", errors="replace")
    reader = csv.reader(io.StringIO(raw))
    rows = list(reader)
    text = "\n".join(", ".join(row) for row in rows)
    return [{"text": text, "metadata": {}}]
