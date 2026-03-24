"""Shared test helpers for creating test documents."""

from pathlib import Path


def create_test_pdf(path: Path, pages: list[str]) -> Path:
    """Create a multi-page PDF with given page texts."""
    import pymupdf

    doc = pymupdf.open()
    for text in pages:
        page = doc.new_page()
        page.insert_text((72, 72), text)
    doc.save(str(path))
    doc.close()
    return path


def create_test_docx(path: Path, paragraphs: list[str]) -> Path:
    """Create a DOCX with given paragraphs."""
    from docx import Document

    doc = Document()
    for para in paragraphs:
        doc.add_paragraph(para)
    doc.save(str(path))
    return path


def create_test_pptx(path: Path, slides: list[str]) -> Path:
    """Create a PPTX with one text box per slide."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for text in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(3))
        txBox.text_frame.text = text
    prs.save(str(path))
    return path
